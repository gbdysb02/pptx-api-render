from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def generate_pptx(companies, output_path="Company_Summary_Deck_Full.pptx"):
    prs = Presentation("template.pptx")  # use uploaded custom template
    layout = prs.slide_layouts[0]  # or choose a specific layout index as needed

    headers = ["Company", "Website", "Overview", "Revenue (PLNm)", "EBIT (PLNm)", "Shareholders"]
    col_widths = [1.5, 1.5, 3.2, 1.2, 1.2, 1.6]

    for i in range(0, len(companies), 4):
        batch = companies[i:i + 4]
        slide = prs.slides.add_slide(layout)

        # Title (can be skipped if the layout already includes one)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_frame = title_shape.text_frame
        title_frame.text = "Company Overview Summary"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        table = slide.shapes.add_table(len(batch) + 1, len(headers), Inches(0.5), Inches(1), Inches(9), Inches(4)).table

        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)

        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True

        for row_idx, company in enumerate(batch, start=1):
            table.cell(row_idx, 0).text = company.get("Company Name", "")
            table.cell(row_idx, 1).text = company.get("Website", "")
            table.cell(row_idx, 2).text = company.get("Overview", "")
            table.cell(row_idx, 3).text = str(company.get("Revenue", ""))
            table.cell(row_idx, 4).text = str(company.get("EBIT", ""))
            table.cell(row_idx, 5).text = company.get("Shareholders", "")

    prs.save(output_path)
    return output_path
